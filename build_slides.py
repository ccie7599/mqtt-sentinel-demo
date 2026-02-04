#!/usr/bin/env python3
"""
Build MQTT Sentinel customer presentation slides.
Simple style: black text, white background, no marketing fluff.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
DARK_BLUE = RGBColor(0x00, 0x3D, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)


def add_slide(title_text, layout_index=6):
    """Add a blank slide and return it."""
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_title(slide, text, top=Inches(0.4), left=Inches(0.7), width=Inches(12), size=Pt(32)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = BLACK
    p.font.bold = True
    return txBox


def add_subtitle(slide, text, top=Inches(1.0), left=Inches(0.7), width=Inches(12), size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = GRAY
    return txBox


def add_body_text(slide, lines, top=Inches(1.8), left=Inches(0.7), width=Inches(11.5), size=Pt(16), line_spacing=Pt(28)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = BLACK
        p.space_after = line_spacing
    return txBox


def add_table(slide, headers, rows, top=Inches(2.0), left=Inches(0.7), col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(12 / n_cols)] * n_cols

    table_width = sum(col_widths)
    row_height = Inches(0.45)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BLACK
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_image(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)


# ─────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Title")
add_title(slide, "MQTT Sentinel", top=Inches(2.5), size=Pt(44))
add_subtitle(slide, "Distributed MQTT Security Platform", top=Inches(3.3), size=Pt(24))

# Thin line
from pptx.util import Emu
shape = slide.shapes.add_shape(
    1,  # rectangle
    Inches(0.7), Inches(3.1), Inches(3), Pt(2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# ─────────────────────────────────────────────────────────────────
# SLIDE 2: Customer Architecture (Current State)
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Customer Arch")
add_title(slide, "Current Architecture")
add_subtitle(slide, "WMS MQTT deployment — current state")

customer_img = os.path.expanduser("~/customer.png")
add_image(slide, customer_img, Inches(0.5), Inches(1.6), width=Inches(12.3))

# ─────────────────────────────────────────────────────────────────
# SLIDE 3: Customer Requirements
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Requirements")
add_title(slide, "Requirements")

add_table(slide, ["Parameter", "Value"], [
    ["Protocol", "MQTT 3.1.1"],
    ["Use Case", "Server-to-device alert notification (one-way, subscribe-only clients)"],
    ["Topic Structure", "Unique device ID per topic (e.g. wyse3299349539363084573)"],
    ["Message Pattern", "Per-topic delivery only, no broadcast"],
    ["Message Rate", "~600 msgs/sec aggregate (spike)"],
    ["Payload", "\"ALERT\" notification (~5 bytes current, up to 100 bytes future)"],
    ["QoS", "Level 1 (at least once delivery)"],
    ["Retention", "3 days"],
    ["Concurrency", "1.5 million concurrent connections"],
    ["Growth", "500K additional connections per year"],
], top=Inches(1.5), col_widths=[Inches(3), Inches(9)])

# ─────────────────────────────────────────────────────────────────
# SLIDE 4: Problem Statement
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Problem")
add_title(slide, "Problem Statement")

add_body_text(slide, [
    "1.  MQTT is a binary protocol — existing WAF infrastructure (Akamai, F5) cannot inspect it.",
    "2.  Origin ingresses are directly exposed to 1.5M+ device connections with no application-layer filtering.",
    "3.  No distributed DDoS protection at the MQTT protocol level — rate limiting, auth brute-force prevention, and packet validation do not exist at the edge.",
    "4.  No visibility into MQTT-specific security events — connection anomalies, malicious payloads, and auth failures are not surfaced in existing monitoring.",
    "5.  Scaling to 1.5M+ connections (growing 500K/year) requires a purpose-built proxy tier, not just broker scaling.",
], top=Inches(1.5), size=Pt(17), line_spacing=Pt(32))

# ─────────────────────────────────────────────────────────────────
# SLIDE 5: Solution Overview
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Solution")
add_title(slide, "Solution Overview")
add_subtitle(slide, "MQTT Sentinel — layered security architecture")

arch_img = os.path.join(SCRIPT_DIR, "images", "mqtt-sentinel-architecture.png")
add_image(slide, arch_img, Inches(0.5), Inches(1.6), width=Inches(12.3))

# ─────────────────────────────────────────────────────────────────
# SLIDE 6: Platform-Level Network Security
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Platform Security")
add_title(slide, "Platform-Level Network Security")
add_subtitle(slide, "Infrastructure protection before traffic reaches MQTT Sentinel")

add_table(slide, ["Protection Layer", "Detail"], [
    ["Linode Region-Level DDoS", "Always-on DDoS mitigation at every Linode region edge. Drops volumetric attacks (UDP/ICMP floods, amplification), protocol attacks (SYN floods), and network anomalies automatically."],
    ["Akamai Prolexic Routed On-Demand", "Network-layer scrubbing for sustained large-scale attacks. Traffic rerouted to Prolexic scrubbing centers, cleaned, and returned via GRE tunnels. Multi-terabit capacity."],
    ["Coverage", "All platform IP space used by MQTT Sentinel proxy and broker instances is protected."],
    ["Result", "Network-layer (L3/L4) volumetric attacks are absorbed at the infrastructure edge before any MQTT Sentinel component processes traffic."],
], top=Inches(1.8), col_widths=[Inches(3.5), Inches(8.5)])

# ─────────────────────────────────────────────────────────────────
# SLIDE 7: Proxy Layer
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Proxy Layer")
add_title(slide, "Proxy Layer — Distributed Edge Security")

proxy_img = os.path.join(SCRIPT_DIR, "images", "proxy-layer-detail.png")
add_image(slide, proxy_img, Inches(0.3), Inches(1.4), width=Inches(12.7))

# ─────────────────────────────────────────────────────────────────
# SLIDE 7: Proxy Layer Details
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Proxy Details")
add_title(slide, "Proxy Layer — Capabilities")

add_table(slide, ["Capability", "Detail"], [
    ["TLS Termination", "TLS 1.2/1.3 offload at the edge, plain MQTT internally"],
    ["Rate Limiting", "6 tiers: global (10K/s), per-IP (100/s), CONNECT (10/s), PUBLISH (100/s), SUBSCRIBE (20/s), per-client (50/s)"],
    ["Authentication", "Auth callout to customer PerconaDB. Fail-closed. Results cached 60s."],
    ["MQTT Validation", "MQTT 3.1.1 packet parsing — protocol level, packet type, size enforcement"],
    ["DDoS Protection", "Distributed multi-region fleet absorbs volumetric attacks. Horizontal scaling by adding nodes."],
    ["Observability", "Prometheus metrics: connections, rate limit events, auth latency, packet counts"],
], top=Inches(1.5), col_widths=[Inches(2.5), Inches(9.5)])

# ─────────────────────────────────────────────────────────────────
# SLIDE 8: Core Broker — Security Inspection
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Inspection")
add_title(slide, "Core Broker — Security Inspection")

inspect_img = os.path.join(SCRIPT_DIR, "images", "security-inspection.png")
add_image(slide, inspect_img, Inches(0.3), Inches(1.4), width=Inches(12.7))

# ─────────────────────────────────────────────────────────────────
# SLIDE 9: Bridge and Origin Protection
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Bridge")
add_title(slide, "Bridge Service — Origin Protection")

bridge_img = os.path.join(SCRIPT_DIR, "images", "bridge-and-origin.png")
add_image(slide, bridge_img, Inches(0.3), Inches(1.4), width=Inches(12.7))

# ─────────────────────────────────────────────────────────────────
# SLIDE 10: Origin Protection Detail
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Origin Protection")
add_title(slide, "Origin Protection — How It Works")

add_body_text(slide, [
    "1.  Bridge Service subscribes to messages on the core broker via MQTT.",
    "2.  Messages are converted from MQTT to WebSocket (WSS) frames.",
    "3.  WebSocket traffic passes through the customer's existing WAF (Akamai / F5) — the WAF can now inspect what was previously opaque binary MQTT.",
    "4.  Clean traffic arrives at customer origin via Envoy, then to Mosquitto.",
    "5.  Origin never receives raw MQTT from the internet. All inbound connections terminate at the distributed proxy layer.",
    "",
    "Result: existing WAF investment is leveraged for MQTT traffic without protocol changes on the origin.",
], top=Inches(1.5), size=Pt(17), line_spacing=Pt(30))

# ─────────────────────────────────────────────────────────────────
# SLIDE 11: Scale and Performance
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Scale")
add_title(slide, "Scale and Performance")

add_table(slide, ["Metric", "Target"], [
    ["Concurrent Connections", "1,500,000+ (scaling 500K/year)"],
    ["Message Throughput", "600+ msgs/sec aggregate"],
    ["Auth Latency (P99)", "< 10ms"],
    ["Message Delivery (P99)", "Regional: < 50ms, Cross-region: < 200ms"],
    ["Availability", "99.99%"],
    ["Message Retention", "72 hours (3 days)"],
    ["QoS", "Level 1 (at least once)"],
    ["Protocol", "MQTT 3.1.1"],
], top=Inches(1.5), col_widths=[Inches(4), Inches(8)])

# ─────────────────────────────────────────────────────────────────
# SLIDE 12: Demo
# ─────────────────────────────────────────────────────────────────
slide = add_slide("Demo")
add_title(slide, "Live Demo", top=Inches(2.5), size=Pt(44))

shape = slide.shapes.add_shape(
    1,
    Inches(0.7), Inches(3.1), Inches(3), Pt(2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_body_text(slide, [
    "1.  Platform health and connection metrics",
    "2.  Device authentication (valid + invalid credentials)",
    "3.  Rate limiting under burst traffic",
    "4.  Security inspection — malicious payload detection",
    "5.  Grafana dashboard — real-time security events",
], top=Inches(3.5), size=Pt(18), line_spacing=Pt(32))

# ─────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────
output_path = os.path.join(SCRIPT_DIR, "mqtt-sentinel-presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
